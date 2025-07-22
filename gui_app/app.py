import streamlit as st
import subprocess
import tempfile
import os
from pptx import Presentation
import whisper

# Extract text from pptx
def extract_text_from_pptx(pptx_file):
    try:
        prs = Presentation(pptx_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        return f"Error extracting text from PPTX: {str(e)}"

# Transcribe video using Whisper
def transcribe_video(video_path):
    try:
        model = whisper.load_model("base")
        result = model.transcribe(video_path)
        return result["text"]
    except Exception as e:
        return f"Error transcribing video: {str(e)}"

# Summarize with Mistral via Ollama
def summarize_with_ollama(input_text):
    try:
        prompt = "Summarize the following combined meeting content including transcript, slides, and video discussion.\n\n" + input_text
        
        # Use stdin to pass the prompt to ollama
        result = subprocess.run(
            ["ollama", "run", "mistral"], 
            input=prompt, 
            capture_output=True, 
            text=True, 
            timeout=300  # 5 minute timeout
        )
        
        if result.returncode != 0:
            return f"Error running Ollama: {result.stderr}"
        
        return result.stdout.strip()
    except subprocess.TimeoutExpired:
        return "Error: Ollama request timed out (5 minutes)"
    except Exception as e:
        return f"Error: {str(e)}"

# Streamlit UI
st.title("📊 Meeting Summarizer (Transcript + PPTX + Video)")

uploaded_transcript = st.file_uploader("Upload Meeting Transcript (.txt)", type=["txt"])
uploaded_pptx = st.file_uploader("Upload Presentation (.pptx)", type=["pptx"])
uploaded_video = st.file_uploader("Upload Meeting Video (.mp4, .mov, etc.)", type=["mp4", "mov", "mkv"])

if st.button("Generate Summary"):
    combined_text = ""
    error_messages = []

    if uploaded_transcript:
        try:
            transcript_text = uploaded_transcript.read().decode("utf-8")
            combined_text += "\n--- Transcript ---\n" + transcript_text
        except Exception as e:
            error_messages.append(f"Error reading transcript: {str(e)}")

    if uploaded_pptx:
        ppt_text = extract_text_from_pptx(uploaded_pptx)
        if ppt_text.startswith("Error"):
            error_messages.append(ppt_text)
        else:
            combined_text += "\n--- Slides ---\n" + ppt_text

    if uploaded_video:
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as tmp_vid:
                tmp_vid.write(uploaded_video.read())
                tmp_vid_path = tmp_vid.name
            
            st.info("Transcribing video... This may take a few minutes.")
            video_text = transcribe_video(tmp_vid_path)
            os.unlink(tmp_vid_path)
            
            if video_text.startswith("Error"):
                error_messages.append(video_text)
            else:
                combined_text += "\n--- Video ---\n" + video_text
        except Exception as e:
            error_messages.append(f"Error processing video: {str(e)}")

    if error_messages:
        for error in error_messages:
            st.error(error)

    if combined_text.strip() == "":
        st.warning("Please upload at least one input.")
    else:
        st.info("Generating summary with Mistral via Ollama...")
        summary = summarize_with_ollama(combined_text)
        
        if summary.startswith("Error"):
            st.error(summary)
        else:
            st.success("✅ Summary Generated")
            st.text_area("📄 Summary Output", summary, height=400)
            st.download_button("Download Summary", summary, file_name="summary.txt")