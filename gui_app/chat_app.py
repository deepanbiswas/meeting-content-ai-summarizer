import streamlit as st
import subprocess
import tempfile
import os
from pptx import Presentation
import whisper

# Extract text from pptx
def extract_text_from_pptx(pptx_file):
    try:
        prs = Presentation(pptx_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        return f"Error extracting text from PPTX: {str(e)}"

# Transcribe video using Whisper
def transcribe_video(video_path):
    try:
        model = whisper.load_model("base")
        result = model.transcribe(video_path)
        return result["text"]
    except Exception as e:
        return f"Error transcribing video: {str(e)}"

# Ask question using Ollama
def ask_ollama(question, context):
    try:
        prompt = f"""Based on the following meeting content, please answer the question:

MEETING CONTENT:
{context}

QUESTION: {question}

Please provide a detailed answer based only on the information provided in the meeting content above. If the information is not available in the content, please say so."""
        
        result = subprocess.run(
            ["ollama", "run", "mistral"], 
            input=prompt, 
            capture_output=True, 
            text=True, 
            timeout=120
        )
        
        if result.returncode != 0:
            return f"Error running Ollama: {result.stderr}"
        
        return result.stdout.strip()
    except subprocess.TimeoutExpired:
        return "Error: Ollama request timed out (2 minutes)"
    except Exception as e:
        return f"Error: {str(e)}"

# Streamlit UI
st.title("🤖 Chat with Meeting Content")
st.markdown("Upload meeting files and ask questions about the content using AI")

# File upload section
st.header("📁 Upload Content")
col1, col2, col3 = st.columns(3)

with col1:
    uploaded_transcript = st.file_uploader("Transcript (.txt)", type=["txt"])
with col2:
    uploaded_pptx = st.file_uploader("Slides (.pptx)", type=["pptx"])
with col3:
    uploaded_video = st.file_uploader("Video (.mp4, .mov)", type=["mp4", "mov", "mkv"])

# Load content
if st.button("Load Content") or any([uploaded_transcript, uploaded_pptx, uploaded_video]):
    combined_text = ""
    error_messages = []

    if uploaded_transcript:
        try:
            transcript_text = uploaded_transcript.read().decode("utf-8")
            combined_text += "\n--- Transcript ---\n" + transcript_text
            st.success("✅ Transcript loaded")
        except Exception as e:
            error_messages.append(f"Error reading transcript: {str(e)}")

    if uploaded_pptx:
        ppt_text = extract_text_from_pptx(uploaded_pptx)
        if ppt_text.startswith("Error"):
            error_messages.append(ppt_text)
        else:
            combined_text += "\n--- Slides ---\n" + ppt_text
            st.success("✅ Slides loaded")

    if uploaded_video:
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as tmp_vid:
                tmp_vid.write(uploaded_video.read())
                tmp_vid_path = tmp_vid.name
            
            with st.spinner("Transcribing video... This may take a few minutes."):
                video_text = transcribe_video(tmp_vid_path)
            os.unlink(tmp_vid_path)
            
            if video_text.startswith("Error"):
                error_messages.append(video_text)
            else:
                combined_text += "\n--- Video ---\n" + video_text
                st.success("✅ Video transcribed and loaded")
        except Exception as e:
            error_messages.append(f"Error processing video: {str(e)}")

    if error_messages:
        for error in error_messages:
            st.error(error)

    if combined_text.strip():
        st.session_state.content = combined_text
        st.success("🎉 Content loaded! You can now ask questions below.")
    else:
        st.warning("Please upload at least one file to continue.")

# Chat interface
if 'content' in st.session_state:
    st.header("💬 Ask Questions")
    
    # Initialize chat history
    if 'chat_history' not in st.session_state:
        st.session_state.chat_history = []
    
    # Quick action buttons
    st.markdown("**Quick Actions:**")
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        if st.button("📋 Summary"):
            question = "Please provide a concise summary of the main points discussed in this meeting."
            st.session_state.current_question = question
    
    with col2:
        if st.button("📝 Topics"):
            question = "What are the main topics or subjects discussed in this meeting?"
            st.session_state.current_question = question
    
    with col3:
        if st.button("👥 Participants"):
            question = "Who are the participants or speakers in this meeting?"
            st.session_state.current_question = question
    
    with col4:
        if st.button("🔑 Key Points"):
            question = "What are the key decisions, action items, or important points from this meeting?"
            st.session_state.current_question = question
    
    # Question input
    question = st.text_input(
        "Your question:", 
        value=st.session_state.get('current_question', ''),
        placeholder="Ask anything about the meeting content..."
    )
    
    if st.button("Ask Question") and question:
        with st.spinner("🤖 Thinking..."):
            answer = ask_ollama(question, st.session_state.content)
        
        # Add to chat history
        st.session_state.chat_history.append({"question": question, "answer": answer})
        
        # Clear current question
        if 'current_question' in st.session_state:
            del st.session_state.current_question
    
    # Display chat history
    if st.session_state.chat_history:
        st.header("💭 Chat History")
        
        for i, chat in enumerate(reversed(st.session_state.chat_history)):
            with st.expander(f"Q: {chat['question'][:60]}{'...' if len(chat['question']) > 60 else ''}", expanded=(i == 0)):
                st.markdown(f"**Question:** {chat['question']}")
                if chat['answer'].startswith("Error"):
                    st.error(chat['answer'])
                else:
                    st.markdown(f"**Answer:** {chat['answer']}")
        
        if st.button("🗑️ Clear Chat History"):
            st.session_state.chat_history = []
            st.rerun()

else:
    st.info("👆 Please upload and load content first to start asking questions.")

# Sidebar with instructions
with st.sidebar:
    st.header("📖 How to Use")
    st.markdown("""
1. **Upload Files**: Add transcript, slides, or video files
2. **Load Content**: Click 'Load Content' to process files
3. **Ask Questions**: Use the text input or quick action buttons
4. **Review Answers**: Check the chat history for all Q&A

**Supported Files:**
- 📄 Transcripts: .txt files
- 📊 Slides: .pptx files  
- 🎥 Videos: .mp4, .mov, .mkv files

**Tips:**
- Be specific in your questions
- Ask about participants, decisions, topics, etc.
- Use quick actions for common queries
""")
    
    st.header("🔧 Requirements")
    st.markdown("""
- Ollama running with Mistral model
- Python packages: whisper, python-pptx
""")
