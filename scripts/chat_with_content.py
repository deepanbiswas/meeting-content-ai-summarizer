#!/usr/bin/env python3
"""
Interactive Chat with Meeting Content
Allows you to ask questions about transcript, slides, and video content using Ollama
"""

import argparse
import os
import subprocess
import tempfile
from pptx import Presentation
import whisper

def extract_text_from_pptx(pptx_path):
    """Extract text from PowerPoint file"""
    try:
        prs = Presentation(pptx_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        return f"Error extracting text from PPTX: {str(e)}"

def transcribe_video(video_path):
    """Transcribe video using Whisper"""
    try:
        model = whisper.load_model("base")
        result = model.transcribe(video_path)
        return result["text"]
    except Exception as e:
        return f"Error transcribing video: {str(e)}"

def ask_ollama(question, context):
    """Ask a question about the content using Ollama"""
    try:
        prompt = f"""Based on the following meeting content, please answer the question:

MEETING CONTENT:
{context}

QUESTION: {question}

Please provide a detailed answer based only on the information provided in the meeting content above. If the information is not available in the content, please say so."""
        
        # Use stdin to pass the prompt to ollama
        result = subprocess.run(
            ["ollama", "run", "mistral"], 
            input=prompt, 
            capture_output=True, 
            text=True, 
            timeout=120  # 2 minute timeout for questions
        )
        
        if result.returncode != 0:
            return f"Error running Ollama: {result.stderr}"
        
        return result.stdout.strip()
    except subprocess.TimeoutExpired:
        return "Error: Ollama request timed out (2 minutes)"
    except Exception as e:
        return f"Error: {str(e)}"

def load_content(args):
    """Load and combine content from transcript, slides, and video"""
    combined_text = ""
    error_messages = []

    if args.transcript:
        try:
            with open(args.transcript, 'r', encoding='utf-8') as f:
                transcript_text = f.read()
            combined_text += "\n--- Transcript ---\n" + transcript_text
            print(f"✅ Loaded transcript: {args.transcript}")
        except Exception as e:
            error_messages.append(f"Error reading transcript: {str(e)}")

    if args.slides:
        ppt_text = extract_text_from_pptx(args.slides)
        if ppt_text.startswith("Error"):
            error_messages.append(ppt_text)
        else:
            combined_text += "\n--- Slides ---\n" + ppt_text
            print(f"✅ Loaded slides: {args.slides}")

    if args.video:
        print("🎥 Transcribing video... This may take a few minutes.")
        video_text = transcribe_video(args.video)
        
        if video_text.startswith("Error"):
            error_messages.append(video_text)
        else:
            combined_text += "\n--- Video ---\n" + video_text
            print(f"✅ Loaded video: {args.video}")

    if error_messages:
        print("❌ Errors encountered:")
        for error in error_messages:
            print(f"   {error}")

    if combined_text.strip() == "":
        print("⚠️  Please provide at least one input file (transcript, slides, or video).")
        return None

    return combined_text

def interactive_chat(content):
    """Start interactive chat session"""
    print("\n" + "="*60)
    print("🤖 INTERACTIVE CHAT WITH MEETING CONTENT")
    print("="*60)
    print("Ask questions about the loaded content.")
    print("Type 'quit', 'exit', or 'q' to end the session.")
    print("Type 'help' for available commands.")
    print("-"*60)

    while True:
        try:
            question = input("\n💬 Your question: ").strip()
            
            if question.lower() in ['quit', 'exit', 'q']:
                print("👋 Goodbye!")
                break
            
            if question.lower() == 'help':
                print("""
📋 Available commands:
- Ask any question about the meeting content
- 'summary' - Get a quick summary of the content
- 'topics' - List main topics discussed
- 'participants' - List meeting participants
- 'quit'/'exit'/'q' - End the chat session
""")
                continue
            
            if question.lower() == 'summary':
                question = "Please provide a concise summary of the main points discussed in this meeting."
            elif question.lower() == 'topics':
                question = "What are the main topics or subjects discussed in this meeting?"
            elif question.lower() == 'participants':
                question = "Who are the participants or speakers in this meeting?"
            
            if not question:
                print("❓ Please enter a question.")
                continue
            
            print("🔄 Thinking...")
            answer = ask_ollama(question, content)
            
            if answer.startswith("Error"):
                print(f"❌ {answer}")
            else:
                print("\n🤖 Answer:")
                print("-" * 40)
                print(answer)
                print("-" * 40)
                
        except KeyboardInterrupt:
            print("\n\n👋 Chat session ended.")
            break
        except Exception as e:
            print(f"❌ Error: {str(e)}")

def main():
    parser = argparse.ArgumentParser(description="Interactive chat with meeting content using Ollama")
    parser.add_argument("--transcript", help="Path to transcript (.txt) file")
    parser.add_argument("--slides", help="Path to PowerPoint (.pptx) file")
    parser.add_argument("--video", help="Path to meeting video (.mp4, .mov, etc.)")
    
    args = parser.parse_args()
    
    if not any([args.transcript, args.slides, args.video]):
        print("❌ Please provide at least one input file (--transcript, --slides, or --video)")
        parser.print_help()
        return
    
    print("🚀 Loading content...")
    content = load_content(args)
    
    if content:
        interactive_chat(content)

if __name__ == "__main__":
    main()
